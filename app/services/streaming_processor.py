"""
Streaming/chunked file processors for large file handling.

BQ-VZ-LARGE-FILES Phase 1 (M2) + Phase 2 (M4).
All processors yield chunks instead of returning complete results.

- StreamingTabularProcessor: CSV/TSV via pandas chunked reader,
  Parquet via pyarrow.ParquetFile.iter_batches(), JSON/JSONL line-buffered.
- StreamingDocumentProcessor: PDF page-by-page (pypdfium2), DOCX paragraph-by-paragraph.

Phase 2 (M4) improvements:
- pdfplumber opened once per document instead of per-page
- Per-page error handling with PyPDF fallback for individual pages
- Table extraction available in fallback path too
"""

from __future__ import annotations

import io
import json
import logging
import os
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, Iterator, List, Optional

import pyarrow as pa
import pyarrow.parquet as pq

from app.config import settings

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Data types yielded by streaming processors
# ---------------------------------------------------------------------------


@dataclass
class TextBlock:
    """A single block of text extracted from a document page/section."""
    page_num: int
    text: str
    tables: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)


# ---------------------------------------------------------------------------
# Security helpers
# ---------------------------------------------------------------------------

ZIP_BOMB_RATIO = 100  # reject if uncompressed/compressed > 100x


def check_zip_bomb(filepath: Path) -> None:
    """Reject archives with suspiciously high compression ratios."""
    if not zipfile.is_zipfile(filepath):
        return
    with zipfile.ZipFile(filepath, "r") as zf:
        compressed = sum(i.compress_size for i in zf.infolist())
        uncompressed = sum(i.file_size for i in zf.infolist())
        if compressed > 0 and uncompressed / compressed > ZIP_BOMB_RATIO:
            raise ValueError(
                f"Zip bomb detected: compression ratio {uncompressed / compressed:.0f}x "
                f"exceeds {ZIP_BOMB_RATIO}x limit"
            )


def check_file_size(filepath: Path) -> None:
    """Reject files exceeding MAX_UPLOAD_SIZE_GB."""
    max_bytes = settings.max_upload_size_gb * 1024 * 1024 * 1024
    actual = filepath.stat().st_size
    if actual > max_bytes:
        raise ValueError(
            f"File size {actual / (1024**3):.1f}GB exceeds "
            f"{settings.max_upload_size_gb}GB limit"
        )


# Encoding fallback chain (R6 — charset detection)
_ENCODING_CHAIN = ["utf-8", "utf-8-sig", "latin-1"]


def _open_text_with_fallback(filepath: Path):
    """Open a text file with encoding fallback chain. Returns file handle."""
    for enc in _ENCODING_CHAIN:
        try:
            fh = open(filepath, "r", encoding=enc)
            # Read a small probe to verify encoding works
            fh.read(4096)
            fh.seek(0)
            return fh
        except (UnicodeDecodeError, ValueError):
            continue

    # Try chardet if available
    try:
        import chardet
        raw = filepath.read_bytes()[:32768]
        detected = chardet.detect(raw)
        enc = detected.get("encoding", "utf-8")
        fh = open(filepath, "r", encoding=enc)
        fh.read(4096)
        fh.seek(0)
        return fh
    except Exception:
        pass

    # Final fallback: utf-8 with replacement
    return open(filepath, "r", encoding="utf-8", errors="replace")


def _align_arrow_table(table: pa.Table, ref_schema: pa.Schema) -> pa.Table:
    """Cast *table* columns to match *ref_schema*.

    Pandas per-chunk dtype inference can drift (e.g. a column is int64 in
    the first chunk but object/string in a later chunk that contains "N/A").
    The ParquetWriter requires a single schema, so every batch must conform
    to the reference established by the first chunk.

    Strategy per column:
    1. Same type → keep as-is.
    2. Safe cast succeeds → use it.
    3. Numeric target + castable source → coerce via ``pd.to_numeric``
       (only truly non-numeric values become null; valid numbers survive).
    4. Everything else → fill with nulls (schema preserved, some data lost).
    """
    import pandas as pd

    arrays = []
    for field in ref_schema:
        if field.name not in table.column_names:
            arrays.append(pa.nulls(table.num_rows, type=field.type))
            continue

        col = table.column(field.name)
        if col.type == field.type:
            arrays.append(col)
            continue

        # Fast path: safe Arrow cast
        try:
            arrays.append(col.cast(field.type, safe=False))
            continue
        except (pa.ArrowInvalid, pa.ArrowNotImplementedError, pa.ArrowTypeError):
            pass

        # Numeric target: coerce through pandas so only bad values become null
        if pa.types.is_integer(field.type) or pa.types.is_floating(field.type):
            try:
                coerced = pd.to_numeric(col.to_pandas(), errors="coerce")
                float_arr = pa.array(coerced, from_pandas=True)
                arrays.append(float_arr.cast(field.type, safe=False))
                logger.debug(
                    "Schema drift: column %r coerced %s → %s via pd.to_numeric",
                    field.name, col.type, field.type,
                )
                continue
            except Exception:
                pass

        # Last resort: null-fill
        logger.warning(
            "Schema drift: column %r type %s cannot cast to %s, "
            "filling %d values with null",
            field.name, col.type, field.type, len(col),
        )
        arrays.append(pa.nulls(len(col), type=field.type))

    return pa.table(arrays, schema=ref_schema)


# ---------------------------------------------------------------------------
# StreamingTabularProcessor (M2)
# ---------------------------------------------------------------------------


class StreamingTabularProcessor:
    """Yields pyarrow.RecordBatch chunks from tabular files.

    Supported formats: CSV, TSV, Parquet, JSON/JSONL.
    Target batch size: settings.streaming_batch_target_rows rows.
    """

    def __init__(self, filepath: Path, file_type: str):
        self.filepath = filepath
        self.file_type = file_type.lower()
        self.batch_target = settings.streaming_batch_target_rows

    def __iter__(self) -> Iterator[pa.RecordBatch]:
        if self.file_type == "parquet":
            yield from self._iter_parquet()
        elif self.file_type in ("csv", "tsv"):
            yield from self._iter_csv()
        elif self.file_type == "json":
            yield from self._iter_json()
        else:
            raise ValueError(f"Unsupported tabular type for streaming: {self.file_type}")

    # -- Parquet ---------------------------------------------------------

    def _iter_parquet(self) -> Iterator[pa.RecordBatch]:
        """Iterate over row groups in a Parquet file."""
        pf = pq.ParquetFile(str(self.filepath))
        for batch in pf.iter_batches(batch_size=self.batch_target):
            yield batch

    # -- CSV / TSV -------------------------------------------------------

    def _iter_csv(self) -> Iterator[pa.RecordBatch]:
        """Chunked CSV/TSV reading via pandas, yielded as Arrow RecordBatch.

        Pandas infers dtypes independently per chunk, which can produce
        inconsistent Arrow schemas (e.g. int64 in chunk 1, string in
        chunk N when a column has mixed types like "123" and "N/A").
        The ParquetWriter requires all batches to share one schema, so
        we anchor on the first chunk's schema and cast subsequent chunks
        to match — using pandas numeric coercion where possible to
        preserve values, falling back to null only for truly non-castable
        entries.
        """
        import pandas as pd

        sep = "\t" if self.file_type == "tsv" else ","
        fh = _open_text_with_fallback(self.filepath)
        try:
            reader = pd.read_csv(
                fh,
                sep=sep,
                chunksize=self.batch_target,
                low_memory=True,
                on_bad_lines="warn",
            )
            ref_schema: Optional[pa.Schema] = None
            for chunk_df in reader:
                table = pa.Table.from_pandas(chunk_df, preserve_index=False)
                if ref_schema is None:
                    ref_schema = table.schema
                elif table.schema != ref_schema:
                    table = _align_arrow_table(table, ref_schema)
                for batch in table.to_batches():
                    yield batch
        finally:
            fh.close()

    # -- JSON / JSONL ----------------------------------------------------

    def _iter_json(self) -> Iterator[pa.RecordBatch]:
        """Line-buffered JSON (JSONL) or array JSON reading."""
        import pandas as pd

        fh = _open_text_with_fallback(self.filepath)
        try:
            first_char = ""
            for ch in fh.read(256):
                if ch.strip():
                    first_char = ch
                    break
            fh.seek(0)

            if first_char == "[":
                # JSON array — stream with ijson (never loads full file)
                import ijson
                fh.close()  # ijson needs binary mode
                buf = []
                with open(self.filepath, 'rb') as bf:
                    for obj in ijson.items(bf, 'item'):
                        if not isinstance(obj, dict):
                            obj = {"value": obj}
                        buf.append(obj)
                        if len(buf) >= self.batch_target:
                            yield pa.RecordBatch.from_pylist(buf)
                            buf.clear()
                    if buf:
                        yield pa.RecordBatch.from_pylist(buf)
                return
            else:
                # JSONL (one object per line) — already streams correctly
                reader = pd.read_json(
                    fh, lines=True, chunksize=self.batch_target
                )
                for chunk_df in reader:
                    table = pa.Table.from_pandas(chunk_df, preserve_index=False)
                    for batch in table.to_batches():
                        yield batch
                fh.close()
        except Exception:
            if not fh.closed:
                fh.close()
            raise


# ---------------------------------------------------------------------------
# StreamingDocumentProcessor (M2)
# ---------------------------------------------------------------------------


class StreamingDocumentProcessor:
    """Yields TextBlock chunks from document files.

    Supported: PDF (pypdfium2 + pdfplumber), DOCX (python-docx).
    """

    def __init__(self, filepath: Path, file_type: str):
        self.filepath = filepath
        self.file_type = file_type.lower()
        self._pypdf_reader = None

    def __iter__(self) -> Iterator[TextBlock]:
        if self.file_type == "pdf":
            yield from self._iter_pdf()
        elif self.file_type in ("docx", "doc"):
            yield from self._iter_docx()
        elif self.file_type in ("pptx", "ppt"):
            yield from self._iter_pptx()
        else:
            # Route native formats (RTF, ICS, VCF, etc.) through lightweight extractors
            from app.services.format_extractors import can_extract, extract_text_blocks
            if can_extract(self.file_type):
                yield from extract_text_blocks(self.filepath, self.file_type)
            else:
                raise ValueError(f"Unsupported document type: {self.file_type}")

    # -- PDF (pypdfium2 + pdfplumber) ------------------------------------

    def _iter_pdf(self) -> Iterator[TextBlock]:
        """Page-by-page PDF extraction using pypdfium2 with pdfplumber table extraction.

        Phase 2 (M4) improvements:
        - pdfplumber is opened once for the whole document (not per-page)
        - Per-page error handling: if pypdfium2 fails on a specific page,
          fall back to PyPDF for that page only
        - Resources are released per-page to keep memory flat
        """
        try:
            import pypdfium2 as pdfium
        except ImportError:
            logger.warning("pypdfium2 not available, falling back to PyPDF")
            yield from self._iter_pdf_fallback()
            return

        # Open pdfplumber once for table extraction across all pages
        plumber_pdf = self._open_pdfplumber()

        pdf = pdfium.PdfDocument(str(self.filepath))
        try:
            num_pages = len(pdf)
            for page_idx in range(num_pages):
                text = ""
                fallback_used = False
                try:
                    page = pdf[page_idx]
                    textpage = page.get_textpage()
                    text = textpage.get_text_bounded()
                    textpage.close()
                    page.close()
                except Exception as e:
                    # Per-page fallback: if pypdfium2 fails on this page,
                    # try PyPDF for just this page
                    logger.warning(
                        "pypdfium2 failed on page %d, falling back to PyPDF: %s",
                        page_idx, e,
                    )
                    text = self._fallback_page_text(page_idx)
                    fallback_used = True

                tables = self._extract_tables_from_page(plumber_pdf, page_idx)

                if text.strip() or tables:
                    meta = {"page_index": page_idx}
                    if fallback_used:
                        meta["fallback"] = True
                    yield TextBlock(
                        page_num=page_idx + 1,
                        text=text,
                        tables=tables,
                        metadata=meta,
                    )
        finally:
            pdf.close()
            if plumber_pdf is not None:
                plumber_pdf.close()

    def _open_pdfplumber(self):
        """Open pdfplumber once for the document. Returns None if unavailable."""
        try:
            import pdfplumber
            return pdfplumber.open(str(self.filepath))
        except Exception as e:
            logger.debug("pdfplumber not available: %s", e)
            return None

    @staticmethod
    def _extract_tables_from_page(plumber_pdf, page_idx: int) -> List[str]:
        """Extract tables from a single page using an already-open pdfplumber document."""
        if plumber_pdf is None:
            return []
        try:
            if page_idx < len(plumber_pdf.pages):
                page = plumber_pdf.pages[page_idx]
                tables = page.extract_tables()
                return [
                    "\n".join("\t".join(str(c or "") for c in row) for row in table)
                    for table in tables
                    if table
                ]
        except Exception as e:
            logger.debug("pdfplumber table extraction failed for page %d: %s", page_idx, e)
        return []

    def _fallback_page_text(self, page_idx: int) -> str:
        """Extract text from a single page using PyPDF as fallback."""
        try:
            if self._pypdf_reader is None:
                try:
                    from PyPDF2 import PdfReader
                except ImportError:
                    from pypdf import PdfReader
                self._pypdf_reader = PdfReader(str(self.filepath))
            if page_idx < len(self._pypdf_reader.pages):
                return self._pypdf_reader.pages[page_idx].extract_text() or ""
        except Exception as e:
            logger.debug("PyPDF fallback failed for page %d: %s", page_idx, e)
        return ""

    def _iter_pdf_fallback(self) -> Iterator[TextBlock]:
        """Fallback PDF reading via PyPDF2/PyPDF (when pypdfium2 is unavailable)."""
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            from pypdf import PdfReader

        plumber_pdf = self._open_pdfplumber()
        try:
            reader = PdfReader(str(self.filepath))
            for page_idx, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                tables = self._extract_tables_from_page(plumber_pdf, page_idx)
                if text.strip() or tables:
                    yield TextBlock(
                        page_num=page_idx + 1,
                        text=text,
                        tables=tables,
                        metadata={"page_index": page_idx, "fallback": True},
                    )
        finally:
            if plumber_pdf is not None:
                plumber_pdf.close()

    # -- DOCX ------------------------------------------------------------

    def _iter_docx(self) -> Iterator[TextBlock]:
        """Paragraph-by-paragraph DOCX extraction."""
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx not available for DOCX streaming")
            return

        doc = Document(str(self.filepath))
        # Group paragraphs into page-like blocks (DOCX has no real pages)
        current_text_parts: List[str] = []
        block_idx = 0
        PARAGRAPHS_PER_BLOCK = 50

        for para in doc.paragraphs:
            if para.text.strip():
                current_text_parts.append(para.text)

            if len(current_text_parts) >= PARAGRAPHS_PER_BLOCK:
                yield TextBlock(
                    page_num=block_idx + 1,
                    text="\n".join(current_text_parts),
                    metadata={"block_index": block_idx},
                )
                current_text_parts = []
                block_idx += 1

        # Remaining paragraphs
        if current_text_parts:
            yield TextBlock(
                page_num=block_idx + 1,
                text="\n".join(current_text_parts),
                metadata={"block_index": block_idx},
            )

    # -- PPTX ------------------------------------------------------------

    def _iter_pptx(self) -> Iterator[TextBlock]:
        """Slide-by-slide PPTX extraction."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not available for PPTX streaming")
            return

        prs = Presentation(str(self.filepath))
        for slide_idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)

            if texts:
                yield TextBlock(
                    page_num=slide_idx + 1,
                    text="\n".join(texts),
                    metadata={"slide_index": slide_idx},
                )
